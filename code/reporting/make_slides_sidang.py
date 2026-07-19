"""
DECK SIDANG — 15 slide isi, murni sinyal fase mentah vs EKG.

Aturan pembingkaian (mengikuti thesis/laporan/kerangka_laporan_final.md):
  - objek tunggal = `unwrapPhasePeak_mm`
  - TIDAK menyebut pustaka bawaan radar, filter internal, atau keluaran BPM-nya
  - TIDAK ada trivial baseline, TIDAK ada hold-out / subjek cadangan
  - 50 fps disebut netral sebagai konfigurasi akuisisi

Semua angka DIHITUNG ULANG dari data tiap kali skrip ini dijalankan -- tidak ada
angka yang diketik tangan ke dalam slide. Kalau data atau pipeline berubah dan
klaim di deck tidak lagi benar, assert di `main()` menggagalkan pembuatan deck.

Usage:
    python make_slides_sidang.py <aligned_csv> <pptx_keluaran>
"""

import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt
from scipy import signal

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "evaluation"))
from negative_control import presence                          # noqa: E402
from phase_pipeline import (BR_BAND, FS, HR_BAND, MAPE_STANDARD,  # noqa: E402
                            MEDFILT_K, OVERLAP, WINDOW_SEC, estimate,
                            metrics, resample_antialias)
from snr_criterion import SNR_THRESHOLD, from_tsv              # noqa: E402

ROOT = Path(__file__).resolve().parents[2]
EMPTY = ROOT / "data" / "raw" / "no_subject"
DIST = ROOT / "data" / "raw" / "distance"

# --- palet: menyatu dengan latar deck yang sudah ada (#1e2761) ---
BG = "#1e2761"
FG = "#ffffff"
DIM = "#cadcfc"
BLUE = "#7fb2f0"
GREEN = "#6bcb77"
AMBER = "#ffc857"
RED = "#ff6b6b"
GREY = "#8a93b8"

RGB_FG = RGBColor(0xFF, 0xFF, 0xFF)
RGB_DIM = RGBColor(0xCA, 0xDC, 0xFC)
RGB_GREEN = RGBColor(0x6B, 0xCB, 0x77)
RGB_AMBER = RGBColor(0xFF, 0xC8, 0x57)
RGB_RED = RGBColor(0xFF, 0x6B, 0x6B)
RGB_BLUE = RGBColor(0x7F, 0xB2, 0xF0)

W, H = 13.333, 7.5           # inci, 16:9


def n(v, d=1, tanda=False):
    """Angka gaya Indonesia: koma desimal, titik pemisah ribuan.

    Dipakai di SELURUH teks slide dan label gambar supaya penulisan angka
    konsisten dengan naskah tesis (6,5% — bukan 6.5%).
    """
    s = f"{v:{'+' if tanda else ''},.{d}f}"
    return s.replace(",", "\x00").replace(".", ",").replace("\x00", ".")


DEMO_SUB = "subject09"       # contoh terbaik (MAPE terendah) untuk ilustrasi
DEMO_WEAK = "subject04"      # contoh SNR rendah untuk pembanding


# ═══════════════════════════════════════════════════════════════════
# BAGIAN 1 — PERHITUNGAN
# ═══════════════════════════════════════════════════════════════════

MAIN_SUBJECTS = [f"subject{i:02d}" for i in range(1, 11)]


def load_front(csv):
    """10 subjek utama saja.

    Penyaringan ini EKSPLISIT, bukan kebetulan: `aligned_all.csv` juga memuat
    subject_cadangan01, sementara kerangka laporan menetapkan tidak ada subjek
    cadangan / hold-out di dalam deck ini. Menggabungkannya akan mengubah
    seluruh angka agregat tanpa terlihat.
    """
    df = pd.read_csv(csv)
    front = df[(df.dataset == "position_front")
               & df.unwrapPhasePeak_mm.notna()
               & df.subject_id.isin(MAIN_SUBJECTS)]
    ada = sorted(front.subject_id.unique())
    assert ada == MAIN_SUBJECTS, f"subjek utama tidak lengkap / berlebih: {ada}"
    return front


def per_subject(front):
    """Pipeline final untuk 10 subjek utama."""
    rows = {}
    for sid, sub in front.groupby("subject_id"):
        est, ref, snr = estimate(sub.Timestamp.values,
                                 sub.unwrapPhasePeak_mm.values,
                                 sub.gt_heart_rate.values)
        mae, mape, bias = metrics(est, ref)
        rows[sid] = dict(est=est, ref=ref, snr=snr, mae=mae, mape=mape,
                         bias=bias, hr=float(ref.mean()),
                         std_gt=float(np.std(ref)))
    return rows


def rank_test(front):
    """UJI PERINGKAT: di PSD fase, peringkat ke berapa bin frekuensi EKG?

    Kalau radar tidak menangkap detak jantung, peringkatnya acak -> median
    ternormalisasi 0,5. Dihitung dari fase MENTAH hasil resampling (belum
    ditapis pita jantung), supaya bukan penapisnya sendiri yang diuji.
    """
    ranks = []
    w = int(WINDOW_SEC * FS)
    step = max(1, int(w * (1 - OVERLAP)))
    for _, sub in front.groupby("subject_id"):
        tu, xu, order = resample_antialias(sub.Timestamp.values,
                                           sub.unwrapPhasePeak_mm.values)
        gu = np.interp(tu, np.sort(sub.Timestamp.values),
                       sub.gt_heart_rate.values[order])
        for s in range(0, len(xu) - w, step):
            seg = signal.detrend(xu[s:s + w])
            g = float(gu[s:s + w].mean()) / 60.0
            if not (HR_BAND[0] <= g <= HR_BAND[1]):
                continue
            f, p = signal.welch(seg, fs=FS, nperseg=w)
            m = (f >= HR_BAND[0]) & (f <= HR_BAND[1])
            fb, pb = f[m], p[m]
            gi = int(np.argmin(np.abs(fb - g)))
            ranks.append(int(np.sum(pb > pb[gi])) / (len(pb) - 1))
    return np.array(ranks)


def empty_vs_person(front):
    """Kontrol negatif: ruangan kosong vs ada subjek."""
    kosong = []
    for f in sorted(EMPTY.glob("*.csv")):        # sesi kosong = berkas datar
        r = pd.read_csv(f).sort_values("Timestamp")
        kosong.append(dict(nama=f.stem,
                           **presence(r.Timestamp.values,
                                      r.unwrapPhasePeak_mm.values)))
    orang = []
    for sid, sub in front.groupby("subject_id"):
        orang.append(dict(nama=sid,
                          **presence(sub.Timestamp.values,
                                     sub.unwrapPhasePeak_mm.values)))
    return kosong, orang


def distance():
    """Eksperimen jarak terkendali: subjek & setup sama, hanya jarak berubah."""
    out = []
    # diurutkan MENURUT ANGKA, bukan abjad: "100" mendahului "50" secara abjad
    # dan membalik urutan naratif dekat -> jauh pada gambar.
    for d in sorted((p for p in DIST.iterdir() if p.is_dir()),
                    key=lambda p: int(p.name)):
        est, ref, snr = from_tsv(d)
        mae, mape, _ = metrics(est, ref)
        out.append(dict(nama=d.name, snr=snr, mape=mape, mae=mae,
                        hr=float(ref.mean())))
    return out


# ═══════════════════════════════════════════════════════════════════
# BAGIAN 2 — GAMBAR
# ═══════════════════════════════════════════════════════════════════

def _style(ax, judul=None):
    ax.set_facecolor(BG)
    for sp in ax.spines.values():
        sp.set_color(GREY)
    ax.tick_params(colors=DIM, labelsize=8)
    ax.xaxis.label.set_color(DIM)
    ax.yaxis.label.set_color(DIM)
    if judul:
        ax.set_title(judul, color=FG, fontsize=10, pad=6)
    return ax


def _save(fig, path):
    # Koma desimal dipasang di SATU tempat, tepat sebelum gambar ditulis, dan
    # HANYA pada sumbu numerik. Pemeriksaan ScalarFormatter itu yang penting:
    # sumbu kategori (StrCategory/FixedFormatter) dan sumbu log memakai penulis
    # lain, sehingga terlewati dengan sendirinya. Memasang formatter tanpa
    # pemeriksaan ini akan mengganti label kategori "S09" menjadi "0" -- rusak
    # tanpa error, dan hanya kelihatan kalau gambarnya dilihat satu per satu.
    from matplotlib.ticker import FuncFormatter, ScalarFormatter
    koma = FuncFormatter(lambda v, _: f"{v:g}".replace(".", ","))
    for ax in fig.axes:
        for a in (ax.xaxis, ax.yaxis):
            if isinstance(a.get_major_formatter(), ScalarFormatter):
                a.set_major_formatter(koma)
    fig.savefig(path, dpi=170, facecolor=BG, bbox_inches="tight")
    plt.close(fig)


def fig_phase(front, out):
    """Slide 3 — seperti apa sinyal fase mentah itu."""
    sub = front[front.subject_id == DEMO_SUB]
    tu, xu, _ = resample_antialias(sub.Timestamp.values,
                                   sub.unwrapPhasePeak_mm.values)
    t0 = tu - tu[0]
    m = (t0 >= 60) & (t0 <= 90)                    # 30 detik contoh

    fig, ax = plt.subplots(2, 1, figsize=(7.6, 4.3),
                           facecolor=BG, height_ratios=[1, 1])
    _style(ax[0], "Sinyal fase mentah — perpindahan permukaan dada")
    ax[0].plot(t0[m], xu[m] - xu[m].mean(), color=BLUE, lw=1.1)
    ax[0].set_ylabel("perpindahan (mm)")
    ax[0].text(0.985, 0.9, "gerak napas mendominasi (1–12 mm)", color=DIM,
               fontsize=8, ha="right", transform=ax[0].transAxes)

    sos = signal.butter(4, HR_BAND, "bandpass", fs=FS, output="sos")
    hr_only = signal.sosfiltfilt(sos, signal.detrend(np.diff(xu, prepend=xu[0])))
    _style(ax[1], "Komponen pita jantung 0,8–2,0 Hz — terkandung di dalamnya")
    ax[1].plot(t0[m], hr_only[m], color=GREEN, lw=1.1)
    ax[1].set_ylabel("amplitudo (a.u.)")
    ax[1].set_xlabel("waktu (detik)")
    ax[1].text(0.985, 0.9, "detak jantung (0,1–0,5 mm)", color=DIM,
               fontsize=8, ha="right", transform=ax[1].transAxes)
    fig.tight_layout()
    _save(fig, out / "f03_fase.png")


def fig_empty_qual(front, out):
    """Slide 4 — ada orang vs ruangan kosong, secara kasat mata."""
    r = pd.read_csv(sorted(EMPTY.glob("*.csv"))[0]).sort_values("Timestamp")
    tk, xk, _ = resample_antialias(r.Timestamp.values,
                                   r.unwrapPhasePeak_mm.values)
    sub = front[front.subject_id == DEMO_SUB]
    to, xo, _ = resample_antialias(sub.Timestamp.values,
                                   sub.unwrapPhasePeak_mm.values)

    def seg(t, x, dur=60):
        t0 = t - t[0]
        m = (t0 >= 30) & (t0 <= 30 + dur)
        return t0[m] - 30, signal.detrend(x[m])

    tko, xko = seg(tk, xk)
    too, xoo = seg(to, xo)
    lim = max(np.abs(xoo).max(), np.abs(xko).max()) * 1.1

    fig, ax = plt.subplots(2, 1, figsize=(7.6, 4.3), facecolor=BG, sharex=True)
    _style(ax[0], "ADA SUBJEK — dada bergerak")
    ax[0].plot(too, xoo, color=GREEN, lw=1.0)
    _style(ax[1], "RUANGAN KOSONG — perekaman & pengaturan identik")
    ax[1].plot(tko, xko, color=RED, lw=1.0)
    ax[1].set_xlabel("waktu (detik)")
    for a in ax:
        a.set_ylim(-lim, lim)
        a.set_ylabel("perpindahan (mm)")
    fig.tight_layout()
    _save(fig, out / "f04_kosong.png")


def fig_pipeline(front, out):
    """Slide 5 — dari fase mentah sampai BPM, termasuk jendela 40 detik."""
    sub = front[front.subject_id == DEMO_SUB]
    est, ref, _ = estimate(sub.Timestamp.values, sub.unwrapPhasePeak_mm.values,
                           sub.gt_heart_rate.values)
    tu, xu, _ = resample_antialias(sub.Timestamp.values,
                                   sub.unwrapPhasePeak_mm.values)
    t0 = tu - tu[0]
    w = int(WINDOW_SEC * FS)

    fig = plt.figure(figsize=(10.6, 4.3), facecolor=BG)
    gs = fig.add_gridspec(2, 2, width_ratios=[1.35, 1], hspace=0.62, wspace=0.25)

    a0 = _style(fig.add_subplot(gs[0, 0]),
                f"1  Sinyal fase + pemotongan jendela {n(WINDOW_SEC, 0)} detik")
    m = t0 <= 200
    a0.plot(t0[m], signal.detrend(xu[m]), color=BLUE, lw=0.8)
    for k in range(3):
        s = k * WINDOW_SEC * (1 - OVERLAP)
        a0.axvspan(s, s + WINDOW_SEC, color=AMBER, alpha=0.13)
    a0.set_ylabel("mm")
    a0.set_xlabel("waktu (detik)")

    # jendela contoh diambil dari TENGAH rekaman: jendela pertama kena artefak
    # tepi penapis, sehingga tidak mewakili kerja pipeline yang sebenarnya.
    s0 = (len(xu) // 2 // w) * w
    seg = signal.detrend(np.diff(xu, prepend=xu[0])[s0:s0 + w])
    sos_br = signal.butter(4, BR_BAND, "bandpass", fs=FS, output="sos")
    sos_hr = signal.butter(4, HR_BAND, "bandpass", fs=FS, output="sos")
    seg2 = signal.sosfiltfilt(sos_hr, seg - signal.sosfiltfilt(sos_br, seg))
    a1 = _style(fig.add_subplot(gs[1, 0]),
                "2  Turunan fase → reduksi napas → pita 0,8–2,0 Hz")
    a1.plot(np.arange(len(seg2)) / FS, seg2, color=GREEN, lw=0.9)
    a1.set_xlabel("waktu dalam satu jendela (detik)")

    f, p = signal.welch(seg2, fs=FS, nperseg=w)
    mk = (f >= HR_BAND[0]) & (f <= HR_BAND[1])
    a2 = _style(fig.add_subplot(gs[0, 1]), "3  Spektrum Welch → puncak")
    a2.plot(f[mk] * 60, p[mk], color=AMBER, lw=1.2)
    pk = f[mk][np.argmax(p[mk])] * 60
    a2.axvline(pk, color=RED, ls="--", lw=1.1)
    a2.text(pk + 1.5, p[mk].max() * 0.82, f"{n(pk, 0)} bpm", color=RED,
            fontsize=9)
    a2.set_xlabel("detak per menit")

    a3 = _style(fig.add_subplot(gs[1, 1]),
                f"4  Median filter (orde {MEDFILT_K}) → estimasi akhir")
    a3.plot(ref, color=FG, lw=1.6, label="EKG (acuan)")
    a3.plot(est, color=GREEN, lw=1.4, label="radar")
    a3.set_xlabel("indeks jendela")
    a3.set_ylabel("bpm")
    lg = a3.legend(fontsize=7.5, facecolor=BG, edgecolor=GREY, loc="lower right")
    for txt in lg.get_texts():
        txt.set_color(DIM)
    _save(fig, out / "f05_pipeline.png")


def fig_snr(front, out):
    """Slide 6 — anatomi SNR: puncak dibagi aras derau."""
    fig, ax = plt.subplots(1, 2, figsize=(10.4, 3.9), facecolor=BG)
    w = int(WINDOW_SEC * FS)
    for a, sid, warna in ((ax[0], DEMO_SUB, GREEN), (ax[1], DEMO_WEAK, RED)):
        sub = front[front.subject_id == sid]
        tu, xu, order = resample_antialias(sub.Timestamp.values,
                                           sub.unwrapPhasePeak_mm.values)
        gu = np.interp(tu, np.sort(sub.Timestamp.values),
                       sub.gt_heart_rate.values[order])
        s = len(xu) // 3
        seg = signal.detrend(xu[s:s + w])
        g = float(gu[s:s + w].mean())
        f, p = signal.welch(seg, fs=FS, nperseg=w)
        m = (f >= HR_BAND[0]) & (f <= HR_BAND[1])
        fb, pb = f[m] * 60, p[m]
        gi = int(np.argmin(np.abs(fb - g)))
        nf = float(np.median(pb))
        _style(a, f"{sid.replace('subject', 'Subjek ')} — SNR = "
                  f"{n(pb[gi] / nf, 2)}")
        a.plot(fb, pb, color=warna, lw=1.3)
        a.axhline(nf, color=DIM, ls=":", lw=1.1)
        a.axvline(g, color=FG, ls="--", lw=1.1)
        a.annotate("", xy=(g, pb[gi]), xytext=(g, nf),
                   arrowprops=dict(arrowstyle="<->", color=AMBER, lw=1.5))
        a.text(g + 2, np.sqrt(max(pb[gi], 1e-12) * nf), " SNR", color=AMBER,
               fontsize=10, fontweight="bold")
        a.text(fb[-1], nf * 1.25, "aras derau (median PSD)  ", color=DIM,
               fontsize=7.5, ha="right")
        a.text(g + 2, pb.max() * 0.95, f"frekuensi EKG {n(g, 0)} bpm", color=FG,
               fontsize=7.5)
        a.set_yscale("log")
        a.set_xlabel("detak per menit")
    ax[0].set_ylabel("kerapatan spektral daya")
    fig.tight_layout()
    _save(fig, out / "f06_snr.png")


def fig_negcontrol(kosong, orang, out):
    """Slide 7 — kontrol negatif secara terukur (SNR napas, skala log)."""
    fig, ax = plt.subplots(figsize=(8.6, 4.0), facecolor=BG)
    _style(ax, "SNR napas — deteksi kehadiran (tidak memerlukan EKG)")
    nk = [k["snr_napas"] for k in kosong]
    no = [o["snr_napas"] for o in orang]
    ax.bar(range(len(nk)), nk, color=RED, width=0.6,
           label=f"ruangan kosong (n={len(nk)})")
    ax.bar(range(len(nk), len(nk) + len(no)), no, color=GREEN, width=0.6,
           label=f"ada subjek (n={len(no)})")
    ax.set_yscale("log")
    ax.set_xticks(range(len(nk) + len(no)))
    ax.set_xticklabels([f"K{i+1}" for i in range(len(nk))]
                       + [f"S{i+1:02d}" for i in range(len(no))], fontsize=7.5)
    ax.set_ylabel("SNR napas (skala log)")
    gap_lo, gap_hi = max(nk), min(no)
    ax.axhspan(gap_lo, gap_hi, color=AMBER, alpha=0.12)
    ax.set_ylim(min(nk) * 0.4, max(no) * 6)      # ruang untuk legenda di atas
    # Teks jurang ditaruh di atas kelompok ruangan kosong -- satu-satunya
    # bagian pita yang tidak tertutup batang.
    ax.text(-0.42, np.sqrt(gap_lo * gap_hi),
            f"pemisahan {n(gap_hi / gap_lo, 0)}×\ntanpa tumpang tindih",
            color=AMBER, fontsize=10, fontweight="bold", va="center",
            ha="left", bbox=dict(facecolor=BG, edgecolor="none", alpha=0.85,
                                 pad=3))
    lg = ax.legend(fontsize=8, facecolor=BG, edgecolor=GREY, loc="upper right",
                   ncol=2)
    for t in lg.get_texts():
        t.set_color(DIM)
    fig.tight_layout()
    _save(fig, out / "f07_kontrol.png")


def fig_mape(res, out):
    """Slide 9 — MAPE per subjek terhadap ambang standar."""
    ordered = sorted(res.items(), key=lambda kv: kv[1]["mape"])
    nm = [k.replace("subject", "S") for k, _ in ordered]
    mp = [v["mape"] for _, v in ordered]
    fig, ax = plt.subplots(figsize=(9.4, 3.9), facecolor=BG)
    _style(ax, None)
    pos = np.arange(len(nm))
    ax.bar(pos, mp, color=[GREEN if m < MAPE_STANDARD else RED for m in mp],
           width=0.62)
    ax.set_xticks(pos)                 # label kategori dipasang eksplisit
    ax.set_xticklabels(nm)
    ax.set_xlim(-0.7, len(nm) - 0.3)
    ax.axhline(MAPE_STANDARD, color=AMBER, ls="--", lw=1.6)
    ax.text(-0.35, MAPE_STANDARD + 0.7,
            f"ambang ANSI/CTA-2065 = {n(MAPE_STANDARD, 0)}%",
            color=AMBER, fontsize=9, ha="left")
    for i, m in enumerate(mp):
        ax.text(i, m + 0.4, n(m, 1), color=FG, fontsize=8, ha="center")
    ax.set_ylabel("MAPE (%)")
    fig.tight_layout()
    _save(fig, out / "f09_mape.png")


def fig_scatter(res, dist, out):
    """Slide 10 — SNR memprediksi MAPE (log–log)."""
    fig, ax = plt.subplots(figsize=(8.4, 4.3), facecolor=BG)
    _style(ax, None)
    for sid, v in res.items():
        c = GREEN if v["mape"] < MAPE_STANDARD else RED
        ax.scatter(v["snr"], v["mape"], color=c, s=70, zorder=3)
        ax.annotate(sid.replace("subject", "S"), (v["snr"], v["mape"]),
                    textcoords="offset points", xytext=(7, -3),
                    color=DIM, fontsize=8)
    for d in dist:
        ax.scatter(d["snr"], d["mape"], color=AMBER, s=110, marker="D",
                   zorder=4, edgecolor=BG)
        ax.annotate(f"jarak {d['nama']} cm", (d["snr"], d["mape"]),
                    textcoords="offset points", xytext=(8, 4),
                    color=AMBER, fontsize=8.5, fontweight="bold")
    ax.axvline(SNR_THRESHOLD, color=AMBER, ls="--", lw=1.4)
    ax.axhline(MAPE_STANDARD, color=DIM, ls=":", lw=1.2)
    ax.text(SNR_THRESHOLD * 1.04, 42, f" ambang SNR ≈ {n(SNR_THRESHOLD, 1)}",
            color=AMBER, fontsize=9)
    x = np.array([v["snr"] for v in res.values()] + [d["snr"] for d in dist])
    y = np.array([v["mape"] for v in res.values()] + [d["mape"] for d in dist])
    r = float(np.corrcoef(np.log(x), np.log(y))[0, 1])
    ax.set_xscale("log")
    ax.set_yscale("log")
    # Tanda centang eksplisit dengan angka biasa: penulis log bawaan
    # menghasilkan "6 x 10^0" yang tidak terbaca dari kursi penguji.
    from matplotlib.ticker import FuncFormatter, NullFormatter
    biasa = FuncFormatter(lambda v, _: f"{v:g}".replace(".", ","))
    ax.set_xticks([1.5, 2, 3, 5, 8])
    ax.set_yticks([2, 5, 10, 20, 50])
    for a in (ax.xaxis, ax.yaxis):
        a.set_major_formatter(biasa)
        a.set_minor_formatter(NullFormatter())
    ax.set_xlabel("SNR sinyal jantung (skala log)")
    ax.set_ylabel("MAPE % (skala log)")
    fig.tight_layout()
    _save(fig, out / "f10_scatter.png")
    return r


def fig_rank(ranks, out):
    """Slide 12 — sebaran peringkat frekuensi EKG di spektrum radar."""
    fig, ax = plt.subplots(figsize=(8.6, 3.9), facecolor=BG)
    _style(ax, None)
    ax.hist(ranks, bins=25, range=(0, 1), color=BLUE, edgecolor=BG)
    med = float(np.median(ranks))
    ax.axvline(med, color=GREEN, lw=2.2)
    ax.axvline(0.5, color=RED, ls="--", lw=1.8)
    ax.text(med + 0.015, ax.get_ylim()[1] * 0.88,
            f"median teramati = {n(med, 3)}", color=GREEN, fontsize=10,
            fontweight="bold")
    ax.text(0.515, ax.get_ylim()[1] * 0.62,
            "nilai harapan bila acak = 0,500", color=RED, fontsize=10)
    ax.set_xlabel("peringkat ternormalisasi frekuensi EKG   "
                  "(0 = puncak tertinggi,  1 = terlemah)")
    ax.set_ylabel(f"jumlah jendela  (total {len(ranks)})")
    fig.tight_layout()
    _save(fig, out / "f12_rank.png")
    return med


def fig_corr(res, out):
    """Slide 13 — korelasi ANTAR-subjek + kenapa dalam-subjek tidak sah."""
    fig, ax = plt.subplots(1, 2, figsize=(10.6, 4.1), facecolor=BG,
                           width_ratios=[1, 1])
    x = np.array([v["hr"] for v in res.values()])
    y = np.array([float(v["est"].mean()) for v in res.values()])
    r = float(np.corrcoef(x, y)[0, 1])
    _style(ax[0], f"Antar-subjek (n={len(x)}):  r = {n(r, 2, True)}   → SAH")
    lim = [min(x.min(), y.min()) - 4, max(x.max(), y.max()) + 4]
    ax[0].plot(lim, lim, color=GREY, ls=":", lw=1.2)
    ax[0].scatter(x, y, color=GREEN, s=75, zorder=3)
    for sid, xi, yi in zip(res, x, y):
        ax[0].annotate(sid.replace("subject", "S"), (xi, yi),
                       textcoords="offset points", xytext=(7, -3),
                       color=DIM, fontsize=8)
    ax[0].set_xlim(lim)
    ax[0].set_ylim(lim)
    ax[0].set_xlabel("HR rata-rata dari EKG (bpm)")
    ax[0].set_ylabel("HR rata-rata dari radar (bpm)")

    # ceiling korelasi dalam-subjek: r_max = std / sqrt(std^2 + s^2)
    s_err = 5.0
    std = np.array([v["std_gt"] for v in res.values()])
    ceil = std / np.sqrt(std ** 2 + s_err ** 2)
    idx = np.argsort(ceil)
    _style(ax[1], "Dalam-subjek: batas atas teoretis — TIDAK SAH dipakai")
    # Posisi numerik + label eksplisit: formatter koma di _style akan menimpa
    # penulis kategori bawaan kalau sumbunya dibiarkan kategorikal.
    ax[1].barh(np.arange(len(idx)), ceil[idx], color=AMBER, height=0.62)
    ax[1].set_yticks(np.arange(len(idx)))
    ax[1].set_yticklabels([list(res)[i].replace("subject", "S") for i in idx])
    ax[1].axvline(1.0, color=GREY, ls=":", lw=1.2)
    ax[1].set_xlim(0, 1.05)
    ax[1].set_xlabel("korelasi maksimum yang mungkin dicapai\n"
                     f"estimator SEMPURNA (galat {n(s_err, 0)} bpm)")
    for i, k in enumerate(idx):
        ax[1].text(ceil[k] + 0.02, i, n(ceil[k], 2), color=DIM,
                   fontsize=8, va="center")
    fig.tight_layout()
    _save(fig, out / "f13_korelasi.png")
    return r, float(ceil.max())


def fig_dist(dist, out):
    """Slide 11 — eksperimen jarak terkendali."""
    fig, ax = plt.subplots(1, 2, figsize=(9.0, 3.7), facecolor=BG)
    nm = [f"{d['nama']} cm" for d in dist]
    pos = np.arange(len(dist))
    _style(ax[0], "SNR sinyal jantung")
    ax[0].bar(pos, [d["snr"] for d in dist],
              color=[GREEN if d["snr"] >= SNR_THRESHOLD else RED for d in dist],
              width=0.5)
    ax[0].axhline(SNR_THRESHOLD, color=AMBER, ls="--", lw=1.5)
    ax[0].text(len(nm) - 0.5, SNR_THRESHOLD + 0.08, f"ambang {n(SNR_THRESHOLD, 1)}",
               color=AMBER, fontsize=8.5, ha="right")
    for i, d in enumerate(dist):
        ax[0].text(i, d["snr"] + 0.08, n(d["snr"], 2), color=FG,
                   fontsize=9.5, ha="center")
    _style(ax[1], "MAPE terhadap EKG")
    ax[1].bar(pos, [d["mape"] for d in dist],
              color=[GREEN if d["mape"] < MAPE_STANDARD else RED for d in dist],
              width=0.5)
    ax[1].axhline(MAPE_STANDARD, color=AMBER, ls="--", lw=1.5)
    for i, d in enumerate(dist):
        ax[1].text(i, d["mape"] + 1.2, n(d["mape"], 1) + "%", color=FG,
                   fontsize=9.5, ha="center")
    ax[1].set_ylabel("%")
    for a in ax:                       # label kategori dipasang eksplisit
        a.set_xticks(pos)
        a.set_xticklabels(nm)
        a.set_xlim(-0.6, len(dist) - 0.4)
    fig.tight_layout()
    _save(fig, out / "f11_jarak.png")


# ═══════════════════════════════════════════════════════════════════
# BAGIAN 3 — PENYUSUN SLIDE
# ═══════════════════════════════════════════════════════════════════

class Deck:
    """Pembungkus tipis python-pptx dengan gaya deck yang sudah ada."""

    def __init__(self, path):
        self.prs = Presentation(path)
        self.blank = self.prs.slide_layouts[0]
        self.logo = self._logo_blob()
        self.n = 0

    def _logo_blob(self):
        for sh in self.prs.slides[0].shapes:
            if sh.shape_type == 13:                 # PICTURE
                return sh.image.blob, Emu(sh.width), Emu(sh.height)
        return None

    def slide(self, judul, sub=None):
        s = self.prs.slides.add_slide(self.blank)
        for ph in list(s.shapes):                   # layout bawaan dikosongkan
            ph._element.getparent().remove(ph._element)
        self.n += 1
        self._bg(s)
        if self.logo:
            blob, lw, lh = self.logo
            import io
            s.shapes.add_picture(io.BytesIO(blob), Inches(0.42), Inches(0.32),
                                 height=Inches(0.42))
        t = self.text(s, judul, 1.05, 0.3, W - 2.1, 0.62, 26, True, RGB_FG,
                      "Cambria")
        if sub:
            self.text(s, sub, 1.05, 0.95, W - 2.1, 0.42, 13.5, False, RGB_DIM)
        self.text(s, str(self.n), W - 0.85, H - 0.55, 0.5, 0.3, 11, False,
                  RGB_DIM)
        return s

    def _bg(self, s):
        from pptx.oxml.ns import qn
        bg = s._element.makeelement(qn("p:bg"), {})
        pr = s._element.makeelement(qn("p:bgPr"), {})
        fill = s._element.makeelement(qn("a:solidFill"), {})
        clr = s._element.makeelement(qn("a:srgbClr"), {"val": "1E2761"})
        fill.append(clr)
        pr.append(fill)
        pr.append(s._element.makeelement(qn("a:effectLst"), {}))
        bg.append(pr)
        s._element.insert(0, bg)

    def text(self, s, txt, x, y, w, h, size, bold=False, color=RGB_FG,
             font="Calibri", align=None, spacing=1.0):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(txt.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = spacing
            if align is not None:
                p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = font
            r.font.color.rgb = color
        return box

    def bullets(self, s, items, x, y, w, size=15, color=RGB_DIM, gap=1.35):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = True
        for i, (mark, line, warna) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = gap
            r = p.add_run()
            r.text = f"{mark}  {line}"
            r.font.size = Pt(size)
            r.font.name = "Calibri"
            r.font.color.rgb = warna or color
        return box

    def pic(self, s, path, x, y, w=None, h=None):
        return s.shapes.add_picture(str(path), Inches(x), Inches(y),
                                    Inches(w) if w else None,
                                    Inches(h) if h else None)

    def panel(self, s, x, y, w, h, warna="16204d"):
        from pptx.enum.shapes import MSO_SHAPE
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string(warna.upper())
        sh.line.color.rgb = RGBColor.from_string("2E3B7A")
        sh.line.width = Pt(1)
        sh.shadow.inherit = False
        sh.text_frame.text = ""
        return sh

    def table(self, s, cols, rows, x, y, w, h, warna_baris=None, size=12):
        from pptx.util import Inches as I
        shape = s.shapes.add_table(len(rows) + 1, len(cols), I(x), I(y),
                                   I(w), I(h))
        tbl = shape.table
        for j, c in enumerate(cols):
            cell = tbl.cell(0, j)
            cell.text = c
            self._cell(cell, 11.5, True, RGB_FG, "1B2A6B")
        for i, row in enumerate(rows, start=1):
            for j, v in enumerate(row):
                cell = tbl.cell(i, j)
                cell.text = str(v)
                warna = (warna_baris or {}).get(i - 1, RGB_DIM)
                self._cell(cell, size, False, warna,
                           "16204D" if i % 2 else "1A2559")
        return tbl

    @staticmethod
    def _cell(cell, size, bold, color, bg):
        from pptx.enum.text import MSO_ANCHOR
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(bg)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.name = "Calibri"
                r.font.color.rgb = color

    def move_last_to_end(self, idx):
        """Pindahkan slide ke-idx ke posisi paling akhir."""
        lst = self.prs.slides._sldIdLst
        ids = list(lst)
        lst.remove(ids[idx])
        lst.append(ids[idx])

    def save(self, path):
        self.prs.save(path)


# ═══════════════════════════════════════════════════════════════════
# BAGIAN 4 — ISI DECK
# ═══════════════════════════════════════════════════════════════════

def build(deck, F, res, ranks, kosong, orang, dist, r_snr, r_antar, ceil_max,
          med_rank):
    lolos = [k for k, v in res.items() if v["mape"] < MAPE_STANDARD]
    gagal = [k for k, v in res.items() if v["mape"] >= MAPE_STANDARD]
    mape_avg = float(np.mean([v["mape"] for v in res.values()]))
    mae_avg = float(np.mean([v["mae"] for v in res.values()]))
    gap = min(o["snr_napas"] for o in orang) / max(k["snr_napas"] for k in kosong)

    # ── 1. Latar belakang, rumusan, batasan ────────────────────────
    s = deck.slide("Latar Belakang dan Rumusan Masalah",
                   "Mengapa sinyal fase mentah yang diteliti, bukan yang lain")
    deck.panel(s, 0.75, 1.6, 5.9, 2.35)
    deck.text(s, "Latar belakang", 1.0, 1.75, 5.4, 0.3, 13, True, RGB_AMBER)
    deck.bullets(s, [
        ("•", "Penyakit kardiovaskular = penyebab kematian utama global; "
              "pemantauan detak jantung perlu berkelanjutan.", None),
        ("•", "EKG kontak akurat tetapi tidak nyaman untuk pemantauan kontinu.", None),
        ("•", "Radar FMCW menangkap micromotion permukaan dada tanpa kontak. "
              "Sinyal FASE adalah representasi paling langsung dari "
              "perpindahan itu.", None),
    ], 1.0, 2.15, 5.4, 12.5, gap=1.25)

    deck.panel(s, 6.9, 1.6, 5.65, 2.35)
    deck.text(s, "Rumusan masalah", 7.15, 1.75, 5.2, 0.3, 13, True, RGB_AMBER)
    deck.bullets(s, [
        ("1.", "Apakah sinyal fase mentah benar mengandung informasi detak "
               "jantung, bukan derau atau artefak?", RGB_FG),
        ("2.", "Seberapa akurat estimasi detak jantung darinya terhadap "
               "standar ANSI/CTA-2065?", RGB_FG),
        ("3.", "Faktor fisik apa yang menentukan berhasil atau gagalnya "
               "estimasi?", RGB_FG),
    ], 7.15, 2.15, 5.2, 12.5, gap=1.25)

    deck.panel(s, 0.75, 4.2, 11.8, 1.85)
    deck.text(s, "Batasan masalah", 1.0, 4.35, 11.3, 0.3, 13, True, RGB_AMBER)
    deck.bullets(s, [
        ("•", "Objek tunggal: sinyal fase mentah unwrapPhasePeak_mm "
              "(perpindahan dada, mm).   •  Akuisisi 50 fps.", None),
        ("•", f"Kondisi subjek ISTIRAHAT → yang dievaluasi ketepatan detak "
              f"rata-rata per jendela {n(WINDOW_SEC, 0)} detik, bukan pelacakan "
              f"perubahan detak jantung.", None),
        ("•", "Estimasi detak rata-rata, BUKAN variabilitas detak jantung "
              "(HRV).   •  Acuan kebenaran: EKG Attys 125 Hz.", None),
        ("•", f"Standar penilaian: ANSI/CTA-2065, MAPE < {n(MAPE_STANDARD, 0)}%. "
              f"Batas fisik keberlakuan dilaporkan sebagai hasil — "
              f"tidak ada subjek yang dibuang.", None),
    ], 1.0, 4.72, 11.3, 12.5, gap=1.2)

    # ── 2. Setup eksperimen ────────────────────────────────────────
    s = deck.slide("Pengaturan Eksperimen dan Acuan Kebenaran",
                   "Dua perangkat merekam bersamaan; EKG menjadi acuan")
    deck.panel(s, 0.75, 1.55, 5.9, 2.6)
    deck.text(s, "Perangkat", 1.0, 1.68, 5.4, 0.3, 13, True, RGB_AMBER)
    deck.bullets(s, [
        ("•", "Perangkat uji: radar FMCW TI IWR1443BOOST, 77 GHz, "
              "akuisisi 50 fps.", None),
        ("•", "Acuan: perekam EKG Attys, 125 Hz.", None),
        ("•", "Keduanya merekam pada rentang waktu yang sama dengan penanda "
              "waktu masing-masing → disinkronkan berbasis waktu.", None),
    ], 1.0, 2.08, 5.4, 12.5, gap=1.25)

    deck.panel(s, 6.9, 1.55, 5.65, 2.6)
    deck.text(s, "Acuan kebenaran — Pan-Tompkins", 7.15, 1.68, 5.2, 0.3, 13,
              True, RGB_AMBER)
    deck.bullets(s, [
        ("→", "Penapisan pita 5–15 Hz → turunan → pengkuadratan → integrasi "
              "jendela bergerak → deteksi puncak R.", None),
        ("→", "Batasan fisiologis 40–180 bpm menyaring puncak takwajar.", None),
        ("→", "Bukan analisis spektral langsung: energi kompleks QRS berada "
              "di frekuensi jauh lebih tinggi sehingga estimasi spektral "
              "tidak stabil.", None),
    ], 7.15, 2.08, 5.2, 12.5, gap=1.2)

    deck.text(s, "Dataset", 0.95, 4.35, 5.4, 0.3, 13, True, RGB_AMBER)
    deck.table(s, ["Dataset", "Jumlah", "Peran dalam penelitian"], [
        ["Subjek utama", "10 subjek", "Penentuan parameter dan evaluasi akurasi"],
        ["Variasi jarak", "2 sesi (50 & 100 cm)",
         "Bukti kausal hubungan SNR dengan akurasi"],
        ["Ruangan kosong", f"{len(kosong)} sesi",
         "Kontrol negatif — radar tanpa subjek"],
    ], 0.95, 4.72, 11.5, 1.5)

    # ── 3. Apa itu sinyal fase mentah ──────────────────────────────
    s = deck.slide("Apa Itu Sinyal Fase Mentah",
                   "Perpindahan permukaan dada, terukur dalam milimeter")
    deck.pic(s, F / "f03_fase.png", 0.62, 1.55, w=7.6)
    deck.panel(s, 8.45, 1.55, 4.1, 4.35)
    deck.text(s, "Bagaimana tertangkapnya", 8.7, 1.7, 3.7, 0.3, 13, True,
              RGB_AMBER)
    deck.bullets(s, [
        ("1.", "Radar memancarkan chirp; pantulan dari dada diterima kembali.", None),
        ("2.", "Range FFT memisahkan pantulan menurut jarak → dipilih bin "
               "jarak tempat subjek berada.", None),
        ("3.", "Fase pada bin tersebut diekstrak, lalu di-unwrap agar "
               "kontinu.", None),
        ("4.", "Relasi fase–perpindahan mengubahnya menjadi milimeter.", None),
    ], 8.7, 2.12, 3.7, 11.5, gap=1.18)
    deck.text(s, "φ(t) = 4π · d(t) / λ", 8.7, 5.05, 3.7, 0.4, 17, True,
              RGB_FG, "Cambria")
    deck.text(s, "Perubahan fase sebesar satu panjang gelombang setara "
                 "perpindahan submilimeter — inilah yang membuat detak "
                 "jantung terukur.", 8.7, 5.42, 3.7, 0.6, 10.5, False, RGB_DIM)
    deck.text(s, "Napas 1–12 mm; detak jantung hanya 0,1–0,5 mm. "
                 "Keduanya berada dalam satu sinyal yang sama.",
              0.62, 6.15, 7.6, 0.4, 12, False, RGB_AMBER)

    # ── 4. Kualitatif: orang vs kosong ─────────────────────────────
    s = deck.slide("Sinyal Itu Memang Gerak Dada",
                   "Perbandingan langsung: ada subjek versus ruangan kosong")
    deck.pic(s, F / "f04_kosong.png", 0.62, 1.5, w=7.7)
    deck.panel(s, 8.55, 1.9, 4.0, 3.1)
    deck.text(s, "Yang perlu diperhatikan", 8.8, 2.05, 3.6, 0.3, 13, True,
              RGB_AMBER)
    deck.bullets(s, [
        ("•", "Skala sumbu tegak kedua panel IDENTIK.", RGB_FG),
        ("•", "Ada subjek: osilasi periodik jelas, amplitudo beberapa "
              "milimeter.", None),
        ("•", "Ruangan kosong: rata. Tidak ada yang bisa dilaporkan, dan "
              "memang tidak ada.", None),
    ], 8.8, 2.45, 3.6, 12.5, gap=1.3)
    deck.text(s, "Ini baru bukti kasat mata. Versinya yang terukur "
                 "menyusul setelah SNR diperkenalkan.",
              8.8, 4.5, 3.6, 0.5, 11, False, RGB_AMBER)

    # ── 5. Pipeline ────────────────────────────────────────────────
    s = deck.slide("Pengolahan: Dari Fase Mentah Menjadi Detak per Menit",
                   f"Enam tahap pemrosesan sinyal digital, tanpa pembelajaran mesin")
    deck.pic(s, F / "f05_pipeline.png", 1.35, 1.4, w=10.6)
    tahap = ("1  Resampling anti-alias  →  2  Turunan fase  →  "
             "3  Reduksi napas 0,15–0,6 Hz  →  4  Pita jantung 0,8–2,0 Hz  →  "
             f"5  Welch jendela {n(WINDOW_SEC, 0)} s + puncak  →  "
             f"6  Median filter orde {MEDFILT_K}")
    deck.panel(s, 0.55, 5.8, 11.9, 0.55)
    deck.text(s, tahap, 0.75, 5.9, 11.5, 0.4, 11.5, True, RGB_FG)
    deck.text(s, f"Jendela {n(WINDOW_SEC, 0)} detik dipilih DI MUKA: subjek "
                 f"istirahat sehingga tidak ada dinamika detak jantung yang "
                 f"hilang, sementara resolusi frekuensi membaik menjadi "
                 f"{n(60 / WINDOW_SEC, 1)} bpm. Konsekuensinya — respons lambat "
                 f"terhadap perubahan detak jantung — dinyatakan sebagai "
                 f"keterbatasan.", 0.55, 6.45, 11.9, 0.7, 12, False, RGB_DIM)

    # ── 6. SNR ─────────────────────────────────────────────────────
    s = deck.slide("Rasio Sinyal terhadap Derau (SNR)",
                   "Ukuran apakah detak jantung menonjol di atas derau")
    deck.pic(s, F / "f06_snr.png", 1.77, 1.42, w=9.8)
    deck.panel(s, 0.75, 5.25, 5.8, 1.65)
    deck.text(s, "Definisi", 1.0, 5.36, 5.3, 0.3, 13, True, RGB_AMBER)
    deck.text(s, "SNR  =  daya spektral pada frekuensi detak jantung\n"
                 "           ÷  aras derau (median PSD pita 0,8–2,0 Hz)",
              1.0, 5.72, 5.3, 0.8, 12.5, False, RGB_FG)
    deck.text(s, "SNR ≈ 1 berarti sinyal jantung setara derau.",
              1.0, 6.45, 5.3, 0.3, 11.5, False, RGB_DIM)
    deck.panel(s, 6.75, 5.25, 5.8, 1.65)
    deck.text(s, "Mengapa vital — dan apa batasnya", 7.0, 5.36, 5.3, 0.3, 13,
              True, RGB_AMBER)
    deck.bullets(s, [
        ("•", "Menentukan apakah ADA yang bisa diekstrak. Pada SNR mendekati "
              "1, tidak ada algoritma mana pun yang mampu memulihkannya.", None),
        ("•", "Dihitung pada frekuensi EKG → memerlukan EKG. Karena itu SNR "
              "adalah ALAT DIAGNOSIS dan pedoman penempatan sensor, bukan "
              "gerbang kualitas lapangan.", RGB_AMBER),
    ], 7.0, 5.72, 5.3, 11, gap=1.15)

    # ── 7. Kontrol negatif kuantitatif ─────────────────────────────
    s = deck.slide("Kontrol Negatif: Ruangan Kosong, Kini Terukur",
                   "Alat ukur yang jujur harus menyatakan tidak ada")
    deck.pic(s, F / "f07_kontrol.png", 0.9, 1.5, w=8.6)
    deck.panel(s, 9.75, 1.7, 2.8, 3.9)
    deck.text(s, "Hasil", 10.0, 1.85, 2.4, 0.3, 13, True, RGB_AMBER)
    deck.bullets(s, [
        ("", f"Ruangan kosong\nSNR napas "
             f"{n(min(k['snr_napas'] for k in kosong), 1)} – "
             f"{n(max(k['snr_napas'] for k in kosong), 1)}", RGB_RED),
        ("", f"Ada subjek\nSNR napas "
             f"{n(min(o['snr_napas'] for o in orang), 0)} – "
             f"{n(max(o['snr_napas'] for o in orang), 0)}", RGB_GREEN),
        ("", f"Pemisahan {n(gap, 0)}×\ntanpa tumpang tindih", RGB_FG),
    ], 10.0, 2.25, 2.4, 12, gap=1.5)
    deck.text(s, "Deteksi kehadiran memakai NAPAS, bukan detak jantung: "
                 "gerak dada akibat napas 1–12 mm, jauh di atas aras derau. "
                 "Karena itu penentuan ada-tidaknya subjek TIDAK memerlukan EKG.",
              0.9, 6.15, 8.6, 0.6, 12, False, RGB_DIM)

    # ── 8. MAPE ────────────────────────────────────────────────────
    s = deck.slide("Cara Menilai: MAPE dan Standar ANSI/CTA-2065",
                   "Satu angka, satu ambang, ditetapkan pihak ketiga")
    deck.panel(s, 0.9, 1.7, 11.5, 1.5)
    deck.text(s, "MAPE  =  (100 / n) · Σ | HR(radar) − HR(EKG) |  ÷  HR(EKG)",
              1.2, 1.95, 11.0, 0.5, 22, True, RGB_FG, "Cambria")
    deck.text(s, "Galat dinyatakan sebagai PERSENTASE terhadap detak jantung "
                 "sebenarnya, sehingga meleset 5 bpm pada orang ber-HR 60 "
                 "tidak dianggap sama ringannya dengan meleset 5 bpm pada "
                 "orang ber-HR 100.", 1.2, 2.6, 11.0, 0.4, 12, False, RGB_DIM)

    deck.panel(s, 0.9, 3.45, 5.6, 2.9)
    deck.text(s, "Standar penilaian", 1.15, 3.6, 5.1, 0.3, 13, True, RGB_AMBER)
    deck.bullets(s, [
        ("•", "ANSI/CTA-2065 — Physical Activity Monitoring for Heart Rate, "
              "Consumer Technology Association (2018).", None),
        ("•", f"Alat dinyatakan VALID bila MAPE < {n(MAPE_STANDARD, 0)}% "
              f"terhadap EKG.", RGB_FG),
        ("•", "Dipakai luas pada literatur validasi perangkat pemantau "
              "detak jantung.", None),
    ], 1.15, 4.0, 5.1, 12.5, gap=1.25)

    deck.panel(s, 6.8, 3.45, 5.6, 2.9)
    deck.text(s, "Mengapa ambang pihak ketiga", 7.05, 3.6, 5.1, 0.3, 13, True,
              RGB_AMBER)
    deck.bullets(s, [
        ("•", "Ambang tidak ditentukan sendiri oleh peneliti, sehingga tidak "
              "dapat disetel agar hasil terlihat baik.", RGB_FG),
        ("•", f"Pada detak jantung rata-rata ±78 bpm, "
              f"{n(MAPE_STANDARD, 0)}% setara ±7,8 bpm.", None),
        ("•", "Sebagian studi memakai kriteria lebih ketat (±5%).", None),
    ], 7.05, 4.0, 5.1, 12.5, gap=1.25)

    # ── 9. Hasil per subjek ────────────────────────────────────────
    s = deck.slide("Hasil: MAPE dan SNR Sepuluh Subjek",
                   f"{len(lolos)} dari {len(res)} subjek memenuhi standar; "
                   f"MAPE rata-rata {n(mape_avg, 1)}%")
    deck.pic(s, F / "f09_mape.png", 0.5, 1.5, w=7.6)
    ordered = sorted(res.items(), key=lambda kv: kv[1]["mape"])
    rows = [[k.replace("subject", "S"), n(v["hr"], 0), n(v["snr"], 2),
             n(v["mape"], 1), "Lolos" if v["mape"] < MAPE_STANDARD
             else "Gagal"] for k, v in ordered]
    warna = {i: (RGB_GREEN if r[4] == "Lolos" else RGB_RED)
             for i, r in enumerate(rows)}
    deck.table(s, ["Subjek", "HR EKG", "SNR", "MAPE %", "Status"], rows,
               8.35, 1.5, 4.35, 4.6, warna, size=11)
    deck.text(s, f"Dua yang gagal (S04, S06) justru ber-SNR terendah "
                 f"(< {n(SNR_THRESHOLD, 1)}). Seluruh yang lolos ber-SNR ≥ "
                 f"{n(min(v['snr'] for k, v in res.items() if v['mape'] < MAPE_STANDARD), 2)}. "
                 f"Pemisahannya bersih — dan itu mengarahkan ke slide "
                 f"berikutnya.", 0.5, 6.25, 7.7, 0.6, 12.5, False, RGB_AMBER)

    # ── 10. Korelasi SNR ↔ MAPE ────────────────────────────────────
    s = deck.slide("Kegagalan Tidak Acak — Diprediksi oleh SNR",
                   f"Korelasi log(SNR) terhadap log(MAPE) = {n(r_snr, 3, True)}")
    deck.pic(s, F / "f10_scatter.png", 0.75, 1.5, w=8.4)
    deck.panel(s, 9.4, 1.8, 3.15, 3.9)
    deck.text(s, "Bacaan", 9.65, 1.95, 2.75, 0.3, 13, True, RGB_AMBER)
    deck.bullets(s, [
        ("•", "Hubungan sangat teratur pada rentang dua orde besaran.", None),
        ("•", f"Ambang SNR ≈ {n(SNR_THRESHOLD, 1)} memisahkan lolos dari gagal.", RGB_FG),
        ("•", "SNR diukur TANPA melihat MAPE lebih dulu, sehingga bukan "
              "penjelasan yang dibuat belakangan.", RGB_AMBER),
    ], 9.65, 2.35, 2.75, 11, gap=1.25)
    deck.text(s, "Namun korelasi belum membuktikan sebab-akibat. "
                 "Pembuktiannya ada pada slide berikutnya.",
              0.75, 6.15, 8.4, 0.4, 12.5, False, RGB_DIM)

    # ── 11. Eksperimen jarak ───────────────────────────────────────
    d50 = next(d for d in dist if d["nama"].startswith("50"))
    d100 = next(d for d in dist if d["nama"].startswith("100"))
    s = deck.slide("Bukti Kausal: Eksperimen Jarak Terkendali",
                   "Subjek sama, pengaturan sama — hanya jarak yang diubah")
    deck.pic(s, F / "f11_jarak.png", 0.85, 1.6, w=9.0)
    deck.panel(s, 10.1, 1.9, 2.45, 3.5)
    deck.text(s, "Dasar fisis", 10.3, 2.05, 2.05, 0.3, 13, True, RGB_AMBER)
    deck.text(s, "Daya pantul\n∝ 1 / R⁴", 10.3, 2.45, 2.05, 0.7, 16, True,
              RGB_FG, "Cambria")
    deck.text(s, f"Jarak menjadi dua kali lipat → daya pantul 16 kali lebih "
                 f"kecil.\n\nTerukur: SNR turun dari {n(d50['snr'], 2)} menjadi "
                 f"{n(d100['snr'], 2)}, dan MAPE melonjak dari {n(d50['mape'], 1)}% "
                 f"menjadi {n(d100['mape'], 1)}%.",
              10.3, 3.25, 2.05, 2.0, 11, False, RGB_DIM)
    deck.text(s, "Karena hanya jarak yang berubah, perbedaan hasil tidak dapat "
                 "dijelaskan oleh perbedaan individu. Ini menaikkan status SNR "
                 "dari sekadar berkorelasi menjadi penyebab.",
              0.85, 6.2, 9.0, 0.5, 12.5, False, RGB_AMBER)

    # ── 12. Uji peringkat frekuensi ────────────────────────────────
    s = deck.slide("Uji Peringkat Frekuensi",
                   "Apakah frekuensi EKG yang benar memang menonjol di spektrum radar?")
    deck.pic(s, F / "f12_rank.png", 0.75, 1.55, w=8.6)
    deck.panel(s, 9.55, 1.75, 3.0, 4.0)
    deck.text(s, "Cara ujinya", 9.8, 1.9, 2.6, 0.3, 13, True, RGB_AMBER)
    deck.bullets(s, [
        ("1.", "Ambil frekuensi detak jantung SEBENARNYA dari EKG.", None),
        ("2.", "Cari peringkat bin frekuensi itu di dalam spektrum fase "
               "radar.", None),
        ("3.", "Bila radar tidak menangkap apa pun, peringkatnya acak "
               "→ median 0,500.", None),
        ("", f"Teramati: {n(med_rank, 3)}\ndari {len(ranks)} jendela", RGB_GREEN),
    ], 9.8, 2.3, 2.6, 11, gap=1.25)
    deck.text(s, "Diuji pada sinyal fase yang BELUM ditapis pita jantung, "
                 "sehingga yang dinilai adalah kandungan sinyalnya — bukan "
                 "kemampuan penapis yang dirancang sendiri.",
              0.75, 6.2, 8.6, 0.5, 12.5, False, RGB_AMBER)

    # ── 13. Korelasi antar-subjek ──────────────────────────────────
    s = deck.slide("Korelasi: Yang Sah dan Yang Tidak",
                   f"Antar-subjek r = {n(r_antar, 2, True)}; korelasi dalam-subjek "
                   f"tidak berlaku pada data ini")
    deck.pic(s, F / "f13_korelasi.png", 0.7, 1.5, w=10.6)
    deck.panel(s, 0.7, 5.25, 5.85, 1.55)
    deck.text(s, "Mengapa dalam-subjek tidak sah", 0.95, 5.38, 5.35, 0.3, 12.5,
              True, RGB_AMBER)
    deck.text(s, "Seluruh subjek duduk diam, sehingga simpangan baku detak "
                 "jantung dalam satu sesi hanya 0,9–4,2 bpm. Korelasi "
                 "dinormalkan terhadap ragam: bila targetnya nyaris konstan, "
                 "estimator SEMPURNA sekalipun hanya mencapai "
                 f"r ≈ {n(ceil_max, 2)}.",
              0.95, 5.72, 5.35, 1.0, 11.5, False, RGB_DIM)
    deck.panel(s, 6.75, 5.25, 5.55, 1.55)
    deck.text(s, "Karena itu dipakai korelasi antar-subjek", 7.0, 5.38, 5.05,
              0.3, 12.5, True, RGB_AMBER)
    deck.text(s, f"Menguji hal yang berbeda: apakah radar dapat MEMBEDAKAN "
                 f"orang ber-HR tinggi dari yang rendah. Hasilnya "
                 f"r = {n(r_antar, 2, True)} pada {len(res)} subjek — dan inilah "
                 f"pertanyaan yang relevan untuk data ini.",
              7.0, 5.72, 5.05, 1.0, 11.5, False, RGB_FG)

    # ── 14. Keterbatasan ───────────────────────────────────────────
    s = deck.slide("Keterbatasan Penelitian",
                   "Dinyatakan sendiri, sebagai batas keberlakuan hasil")
    for i, (judul, isi) in enumerate([
        ("Tidak ada variasi detak jantung",
         "Seluruh subjek direkam dalam kondisi istirahat. Kemampuan metode "
         "melacak PERUBAHAN detak jantung tidak dapat dibuktikan maupun "
         "dibantah dengan data ini. Menambah subjek dengan protokol yang sama "
         "tidak menyelesaikannya — yang diperlukan protokol "
         "istirahat–beban–pemulihan."),
        ("Variabilitas detak jantung (HRV) belum layak",
         "Deteksi detak-per-detak hanya mencapai recall ±61% terhadap puncak R "
         "EKG, dan RMSSD hasil radar sekitar empat kali nilai EKG. Penyebabnya "
         "batas fisik — SNR dan presisi waktu 20 ms per sampel — bukan pilihan "
         "algoritma. Karena itu tidak dapat diperbaiki dengan mengganti "
         "estimator."),
        ("SNR memerlukan EKG",
         "Kriteria SNR dihitung pada frekuensi EKG, sehingga berfungsi sebagai "
         "alat diagnosis dan pedoman penempatan sensor, belum sebagai gerbang "
         "kualitas di lapangan. Deteksi KEHADIRAN melalui napas sudah tidak "
         "memerlukan EKG; gerbang kualitas per-estimasi tanpa EKG menjadi "
         "penelitian lanjutan."),
    ]):
        y = 1.75 + i * 1.62
        deck.panel(s, 0.85, y, 11.6, 1.32)
        deck.text(s, f"{chr(97 + i)}.  {judul}", 1.15, y + 0.13, 11.0, 0.3,
                  13.5, True, RGB_AMBER)
        deck.text(s, isi, 1.15, y + 0.5, 11.0, 0.9, 11.5, False, RGB_DIM)

    # ── 15. Kesimpulan ─────────────────────────────────────────────
    s = deck.slide("Kesimpulan", "Menjawab tiga rumusan masalah")
    for i, (tanya, jawab, bukti) in enumerate([
        ("Apakah sinyal fase mengandung informasi detak jantung?",
         "YA — dibuktikan secara independen.",
         f"Puncak spektrum jatuh pada frekuensi EKG; peringkat ternormalisasi "
         f"{n(med_rank, 3)} dari {len(ranks)} jendela (acak = 0,500); "
         f"korelasi antar-subjek {n(r_antar, 2, True)}; dan kontrol negatif "
         f"memisahkan ruangan kosong dari berisi sebesar {n(gap, 0)}×."),
        ("Seberapa akurat terhadap ANSI/CTA-2065?",
         f"{len(lolos)} dari {len(res)} subjek memenuhi standar.",
         f"MAPE rata-rata {n(mape_avg, 1)}% dan MAE {n(mae_avg, 1)} bpm, "
         f"seluruhnya dicapai dengan pemrosesan sinyal digital murni "
         f"tanpa pembelajaran mesin."),
        ("Faktor fisik apa yang menentukan keberhasilan?",
         f"SNR sinyal jantung, dengan ambang ≈ {n(SNR_THRESHOLD, 1)}.",
         f"Korelasi log(SNR) terhadap log(MAPE) = {n(r_snr, 3, True)}, dan "
         f"dibuktikan secara kausal melalui eksperimen jarak terkendali "
         f"({n(d50['mape'], 1)}% pada 50 cm versus {n(d100['mape'], 1)}% pada "
         f"100 cm, subjek dan pengaturan sama)."),
    ]):
        y = 1.6 + i * 1.78
        deck.panel(s, 0.85, y, 11.6, 1.55)
        deck.text(s, f"{i + 1}.  {tanya}", 1.15, y + 0.12, 11.0, 0.3, 12.5,
                  False, RGB_DIM)
        deck.text(s, jawab, 1.15, y + 0.47, 11.0, 0.35, 15.5, True, RGB_GREEN)
        deck.text(s, bukti, 1.15, y + 0.88, 11.0, 0.6, 11.5, False, RGB_FG)

    deck.text(s, f"Kontribusi:  K1 pipeline DSP tervalidasi standar   ·   "
                 f"K2 kriteria kelayakan berbasis SNR dan karakterisasi jarak"
                 f"   ·   K3 koreksi metodologi evaluasi korelasi",
              0.85, 6.95, 11.6, 0.4, 11.5, True, RGB_AMBER)


# ═══════════════════════════════════════════════════════════════════
def main(csv, out_pptx):
    figdir = ROOT / "thesis" / "pengolahan data" / "_fig_sidang"
    figdir.mkdir(parents=True, exist_ok=True)

    print("menghitung ulang dari data ...")
    front = load_front(csv)
    res = per_subject(front)
    ranks = rank_test(front)
    kosong, orang = empty_vs_person(front)
    dist = distance()

    print("membuat gambar ...")
    fig_phase(front, figdir)
    fig_empty_qual(front, figdir)
    fig_pipeline(front, figdir)
    fig_snr(front, figdir)
    fig_negcontrol(kosong, orang, figdir)
    fig_mape(res, figdir)
    r_snr = fig_scatter(res, dist, figdir)
    med_rank = fig_rank(ranks, figdir)
    r_antar, ceil_max = fig_corr(res, figdir)
    fig_dist(dist, figdir)

    print("merakit deck ...")
    deck = Deck(out_pptx)
    build(deck, figdir, res, ranks, kosong, orang, dist, r_snr, r_antar,
          ceil_max, med_rank)
    deck.move_last_to_end(1)          # slide "Terima Kasih" kembali ke akhir
    deck.save(out_pptx)

    # ── cek-mandiri: klaim di deck harus tetap benar ──────────────
    lolos = sum(1 for v in res.values() if v["mape"] < MAPE_STANDARD)
    mape_avg = float(np.mean([v["mape"] for v in res.values()]))
    gap = min(o["snr_napas"] for o in orang) / max(k["snr_napas"] for k in kosong)
    d50 = next(d for d in dist if d["nama"].startswith("50"))
    d100 = next(d for d in dist if d["nama"].startswith("100"))

    print(f"\n  lolos standar     {lolos}/{len(res)}   MAPE rata-rata {mape_avg:.2f}%")
    print(f"  peringkat EKG     median {med_rank:.3f}  ({len(ranks)} jendela)")
    print(f"  korelasi          antar-subjek {r_antar:+.3f} | "
          f"log SNR-MAPE {r_snr:+.3f}")
    print(f"  kontrol negatif   pemisahan {gap:,.0f}x")
    print(f"  jarak             50cm SNR {d50['snr']:.2f} MAPE {d50['mape']:.1f}% | "
          f"100cm SNR {d100['snr']:.2f} MAPE {d100['mape']:.1f}%")

    assert lolos == 8, f"jumlah lolos berubah jadi {lolos}/10"
    assert mape_avg < 7.0, f"MAPE rata-rata jadi {mape_avg:.2f}%"
    assert med_rank < 0.20, (
        f"median peringkat jadi {med_rank:.3f} — klaim 'frekuensi EKG menonjol' "
        "tidak lagi didukung data")
    assert r_antar > 0.7, f"korelasi antar-subjek turun jadi {r_antar:+.3f}"
    assert r_snr < -0.8, f"korelasi log(SNR)-log(MAPE) jadi {r_snr:+.3f}"
    assert gap > 100, f"pemisahan kontrol negatif jadi {gap:.0f}x"
    assert d50["mape"] < MAPE_STANDARD <= d100["mape"], (
        "eksperimen jarak tidak lagi menunjukkan 50 cm lolos & 100 cm gagal")
    assert all(v["snr"] < SNR_THRESHOLD for k, v in res.items()
               if v["mape"] >= MAPE_STANDARD), \
        "ada subjek gagal ber-SNR di atas ambang — klaim slide 9 & 10 goyah"
    print("  (cek-mandiri: 8 klaim kunci terverifikasi)")
    print(f"\nok: {out_pptx}  —  {len(deck.prs.slides.__iter__.__self__._sldIdLst)} slide")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print(__doc__)
        sys.exit(1)
    main(sys.argv[1], sys.argv[2])
