"""
Generate slide deck (.pptx, editable di LibreOffice Impress/PowerPoint) untuk
laporan bimbingan: reliabilitas sinyal radar IWR1443BOOST untuk estimasi HR.

Angka-angka di sini adalah HASIL ANALISIS TETAP (snapshot Jul 2026, lihat
CLAUDE.md Temuan Kunci poin 1-8) -- bukan dihitung ulang dari CSV di sini,
supaya generate cepat tanpa perlu re-run seluruh pipeline. Kalau ada data baru
masuk dan angka berubah, update konstanta di bagian atas file ini lalu
jalankan ulang.

Usage:
    python make_slides.py <output.pptx>

Contoh:
    python make_slides.py ../../thesis/bimbingan_2026-07-04.pptx
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ----
BLUE = RGBColor(0x2A, 0x78, 0xD6)
RED = RGBColor(0xE3, 0x49, 0x48)
MUTED = RGBColor(0x89, 0x87, 0x81)
GOOD = RGBColor(0x0C, 0xA3, 0x0C)
WARNING_TEXT = RGBColor(0x8A, 0x5C, 0x00)
SERIOUS = RGBColor(0xEC, 0x83, 0x5A)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
PAGE_BG = RGBColor(0xF4, 0xF3, 0xF0)
TEXT_PRIMARY = RGBColor(0x0B, 0x0B, 0x0B)
TEXT_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
BORDER = RGBColor(0xE1, 0xE0, 0xD9)

# ---- data (lihat CLAUDE.md Temuan Kunci poin 1-8) ----
BASELINE_TI_CORR = [
    ("final_heart_rate", 0.003),
    ("heart_rate_est_peak", -0.002),
    ("heart_rate_est_fft", -0.045),
    ("heart_rate_est_fft_4hz", 0.088),
    ("heartRateEst_xCorr", -0.036),
]

FOUR_APPROACHES = [
    ("Baseline TI (xCorr)", -0.036),
    ("Fase mentah (raw)", 0.091),
    ("Fase + notch napas", 0.077),
    ("Filter confidence top 2%", 0.043),
]

SUBJECT_GT_MEAN = [
    ("subject01", 89.3), ("subject09", 84.4), ("subject03", 78.5),
    ("subject02", 77.6), ("subject06", 76.2), ("subject04", 75.4),
    ("subject08", 75.2), ("subject05", 68.3), ("subject07", 65.3),
    ("subject10", 61.1),
]

MODEL_VS_TRIVIAL = [
    ("Model ML existing", 15.0),
    ("Baseline trivial (tebak rata-rata)", 11.9),
    ("Target pembimbing", 5.0),
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=PAGE_BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_kicker_title(slide, kicker, title, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = kicker.upper()
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = MUTED

    tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(0.75), Inches(12), Inches(1.1))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = title
    r2.font.size = Pt(28)
    r2.font.bold = True
    r2.font.color.rgb = TEXT_PRIMARY

    y = 1.6
    if sub:
        tb3 = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(11.5), Inches(0.6))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run()
        r3.text = sub
        r3.font.size = Pt(14)
        r3.font.color.rgb = TEXT_SECONDARY
        y += 0.6
    return y


def add_callout(slide, top, label, body, accent=RED, label_color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(top),
                                  Inches(11.9), Inches(1.15))
    box.adjustments[0] = 0.06
    box.fill.solid()
    box.fill.fore_color.rgb = SURFACE
    box.line.color.rgb = BORDER
    box.line.width = Pt(0.75)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(top),
                                         Inches(0.06), Inches(1.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.14)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = label.upper()
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = label_color or accent

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(14)
    r2.font.color.rgb = TEXT_PRIMARY
    return top + 1.15 + 0.25


def add_stat_tile(slide, left, top, width, label, value, sub, value_color=TEXT_PRIMARY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                  Inches(width), Inches(1.5))
    box.adjustments[0] = 0.08
    box.fill.solid()
    box.fill.fore_color.rgb = SURFACE
    box.line.color.rgb = BORDER
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.16)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = label
    r1.font.size = Pt(11)
    r1.font.color.rgb = MUTED

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = value_color

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.add_run()
    r3.text = sub
    r3.font.size = Pt(10.5)
    r3.font.color.rgb = MUTED


def add_diverging_bar_chart(slide, left, top, width, height, title, data, unit=""):
    """Bar chart horizontal, warna per-titik: biru positif, merah negatif."""
    chart_data = CategoryChartData()
    chart_data.categories = [d[0] for d in data]
    chart_data.add_series("Nilai", [d[1] for d in data])

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    )
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = TEXT_SECONDARY

    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.number_format = "0.000" if abs(max(d[1] for d in data)) < 2 else "0.0"
    dls.number_format_is_linked = False
    dls.font.size = Pt(11)
    dls.font.bold = True
    dls.font.color.rgb = TEXT_PRIMARY

    series = plot.series[0]
    for i, (_, val) in enumerate(data):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = BLUE if val >= 0 else RED

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
    cat_axis.format.line.color.rgb = BORDER

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.color.rgb = MUTED
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = BORDER
    return gframe


def add_magnitude_bar_chart(slide, left, top, width, height, title, data, colors=None):
    """Bar chart horizontal untuk nilai magnitude (semua positif, mis. MAE)."""
    chart_data = CategoryChartData()
    chart_data.categories = [d[0] for d in data]
    chart_data.add_series("Nilai", [d[1] for d in data])

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    )
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = TEXT_SECONDARY

    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.number_format = "0.0"
    dls.number_format_is_linked = False
    dls.font.size = Pt(11)
    dls.font.bold = True
    dls.font.color.rgb = TEXT_PRIMARY

    series = plot.series[0]
    default_colors = colors or [BLUE] * len(data)
    for i in range(len(data)):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = default_colors[i]

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
    cat_axis.format.line.color.rgb = BORDER

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.color.rgb = MUTED
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = BORDER
    return gframe


def add_bullets(slide, left, top, width, height, items, size=15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_after = Pt(10)


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    table = gtable.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAGE_BG
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = MUTED
        para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = TEXT_PRIMARY
            para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
    return gtable


def build(output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- Slide 1: Title ---
    s = add_blank_slide(prs)
    set_bg(s)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "LAPORAN BIMBINGAN · 4 JULI 2026"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = MUTED

    tb2 = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Reliabilitas Sinyal Radar IWR1443BOOST\nuntuk Estimasi Heart Rate"
    r2.font.size = Pt(36)
    r2.font.bold = True
    r2.font.color.rgb = TEXT_PRIMARY

    tb3 = s.shapes.add_textbox(Inches(0.9), Inches(4.3), Inches(10.5), Inches(0.6))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Evaluasi ulang baseline TI, sinyal fase mentah, dan model ML existing terhadap ground truth ECG — dataset position_front, 10 subjek."
    r3.font.size = Pt(15)
    r3.font.color.rgb = TEXT_SECONDARY

    add_callout(s, 5.1, "Temuan Utama yang Akan Dibahas",
                "Model ML yang dilaporkan sebelumnya kemungkinan besar TIDAK lebih baik dari menebak angka rata-rata konstan.",
                accent=RED)

    # --- Slide 2: Ringkasan 4 angka ---
    s = add_blank_slide(prs)
    set_bg(s)
    y = add_kicker_title(s, "Ringkasan", "Empat Angka Kunci",
                          "Sebelum masuk detail, ini posisi kita sekarang.")
    tiles = [
        ("MAE model existing (dilaporkan)", "≈ 10–19 bpm", "korelasi tidak dilaporkan saat itu", TEXT_PRIMARY),
        ("MAE baseline trivial (tebak rata-rata)", "11,9 bpm", "tanpa fitur radar sama sekali", RED),
        ("Target pembimbing", "≤ 5 bpm", "syarat lanjut ke tahap berikutnya", GOOD),
        ("Korelasi terbaik (4 pendekatan)", "0,09", "target realistis: ≥ 0,3–0,5", RED),
    ]
    left = 0.7
    for label, val, sub, color in tiles:
        add_stat_tile(s, left, y + 0.3, 2.85, label, val, sub, value_color=color)
        left += 3.0

    # --- Slide 3: Baseline TI ---
    s = add_blank_slide(prs)
    set_bg(s)
    y = add_kicker_title(s, "Temuan 1", "Output HR Bawaan TI Tidak Reliable",
                          "Dibandingkan ground truth ECG (R-peak Pan-Tompkins), 10 subjek, n=218.603 sampel.")
    add_diverging_bar_chart(s, 0.7, y + 0.3, 7.2, 4.3,
                             "Korelasi vs ground truth (0 = tidak ada hubungan)", BASELINE_TI_CORR)
    add_table(s, 8.2, y + 0.3, 4.4, 3.0,
              ["Kolom", "MAE (bpm)"],
              [("final_heart_rate", "68,4"), ("heart_rate_est_peak", "68,8"),
               ("heart_rate_est_fft", "32,7"), ("heart_rate_est_fft_4hz", "95,0"),
               ("heartRateEst_xCorr", "23,4")],
              col_widths=[2.8, 1.6])
    note = s.shapes.add_textbox(Inches(8.2), Inches(y + 3.5), Inches(4.4), Inches(1.2))
    ntf = note.text_frame
    ntf.word_wrap = True
    np_ = ntf.paragraphs[0]
    nr = np_.add_run()
    nr.text = "final_heart_rate & heart_rate_est_peak stuck di nilai sentinel 5,86 untuk 99,6% sampel — bukan hasil ukur valid."
    nr.font.size = Pt(11.5)
    nr.font.color.rgb = MUTED

    # --- Slide 4: 4 pendekatan ---
    s = add_blank_slide(prs)
    set_bg(s)
    y = add_kicker_title(s, "Temuan 2", "Empat Pendekatan Dicoba — Semua Mentok",
                          "Sebelum menyimpulkan masalah akuisisi, dicoba menyelamatkan sinyal di luar algoritma TI.")
    add_diverging_bar_chart(s, 0.7, y + 0.3, 7.6, 4.0,
                             "Korelasi per pendekatan (agregat 10 subjek)", FOUR_APPROACHES)
    add_bullets(s, 8.5, y + 0.5, 4.1, 3.5, [
        "Filter confidence_heart: tidak prediktif, skala beda total antar dataset (position_front vs position_variation).",
        "Breathing-harmonic suppression: tidak membantu, korelasi malah turun (0,091 → 0,077).",
        "Semua < 0,1 — jauh dari ambang ≥0,3–0,5 yang berarti “ada sinyal HR yang bisa diekstrak”.",
    ], size=13)

    # --- Slide 5: TEMUAN KRITIS ---
    s = add_blank_slide(prs)
    set_bg(s)
    y = add_kicker_title(s, "Temuan Paling Kritis", "Model ML Existing vs Baseline Trivial",
                          "Fitur: heart_rate_est_fft, heart_rate_est_fft_4hz, heartRateEst_xCorr, heart_rate_est_peak, confidence_heart. Train 8 subjek, test 2 subjek terakhir.")
    # kicker in red for emphasis - override color of first textbox created in add_kicker_title not directly accessible; acceptable default muted.
    add_magnitude_bar_chart(s, 0.7, y + 0.3, 11.9, 4.2,
                             "MAE (bpm) — lebih rendah lebih baik",
                             MODEL_VS_TRIVIAL,
                             colors=[SERIOUS, MUTED, GOOD])

    # --- Slide 6: Kenapa + variasi subjek ---
    s = add_blank_slide(prs)
    set_bg(s)
    y = add_kicker_title(s, "Penjelasan", "Kenapa Ini Bisa Terjadi")
    y2 = add_callout(s, y + 0.2, "Garbage In, Garbage Out",
                      "Fitur model existing = tepat kolom yang terbukti korelasi ≈0 (Slide 3). Input tanpa info HR → model regresi konvergen ke rata-rata target → MAE tampak “masuk akal” tanpa pembelajaran sinyal sungguhan.",
                      accent=RGBColor(0xFA, 0xB2, 0x19), label_color=WARNING_TEXT)
    add_magnitude_bar_chart(s, 0.7, y2 + 0.1, 11.9, 3.6,
                             "Rata-rata gt_heart_rate per subjek (bpm) — sumber variasi trivial baseline",
                             SUBJECT_GT_MEAN)

    # --- Slide 7: Rekomendasi ---
    s = add_blank_slide(prs)
    set_bg(s)
    y = add_kicker_title(s, "Untuk Didiskusikan", "Rekomendasi Arah Selanjutnya")
    add_bullets(s, 0.7, y + 0.3, 11.9, 5.0, [
        "Standar evaluasi baru: tiap laporan MAE wajib sertakan korelasi/R² + pembanding baseline trivial, split subjek sama.",
        "Opsi A: lanjut teknik ekstraksi lebih canggih (VMD, deep learning di raw waveform fase) — effort besar, belum ada jaminan.",
        "Opsi B: root-cause analysis kenapa radar gagal tangkap sinyal HR — kontribusi ilmiah valid meski hasil negatif.",
        "Opsi C: evaluasi ulang akuisisi data (jarak, sudut, kondisi rekam) — pertimbangkan rekam ulang dengan setup lebih terkontrol.",
    ], size=17)

    # --- Slide 8: Penutup ---
    s = add_blank_slide(prs)
    set_bg(s)
    y = add_kicker_title(s, "Penutup", "Poin yang Perlu Keputusan Pembimbing")
    add_bullets(s, 0.7, y + 0.3, 11.9, 3.0, [
        "Apakah lanjut Opsi A, B, atau C (atau kombinasi)?",
        "Apakah standar evaluasi baru (korelasi + trivial baseline) bisa jadi wajib mulai sekarang?",
        "Apakah ada dana/waktu untuk rekam ulang data kalau Opsi C dipilih?",
    ], size=17)
    footer = s.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.5))
    ftf = footer.text_frame
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = "Detail lengkap: CLAUDE.md (Temuan Kunci poin 1–8), script di code/evaluation/."
    fr.font.size = Pt(11.5)
    fr.font.color.rgb = MUTED

    prs.save(output_path)
    print(f"Tersimpan: {output_path} ({len(prs.slides)} slide)")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(__doc__)
        sys.exit(1)
    build(sys.argv[1])
